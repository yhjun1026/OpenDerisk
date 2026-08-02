"""Built-in file extractors (RFC 004 §6).

Each extractor is registered via the `@extractor` decorator with one or more
mime glob patterns. Plain-text extractors (txt/md/pdf/docx/pptx) need no
model; image/audio extractors call the `model_caller` supplied by the
orchestrator.

Heavy deps (pdfplumber, python-docx, python-pptx) are imported lazily inside
the extract methods so the module imports cleanly even when optional deps
are missing — the extractor just raises a clear error at extract time.
"""

from __future__ import annotations

import asyncio
import base64
import logging
import mimetypes
import os
import re
import shutil
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Optional

from derisk.knowledge.types import ExtractMode

from . import Extractor, ModelCaller, VerbatimSpec, extractor

logger = logging.getLogger(__name__)


def _file_mtime_iso(path: Path) -> tuple[Optional[str], Optional[int]]:
    try:
        st = path.stat()
        return (
            datetime.fromtimestamp(st.st_mtime, tz=timezone.utc).isoformat(),
            int(st.st_mtime),
        )
    except OSError:
        return None, None


# ---------------------------------------------------------------------------
# Text-family extractors
# ---------------------------------------------------------------------------


@extractor("text", ["text/*", "application/markdown", "application/json", "application/x-yaml", "application/xml"])
class TextExtractor(Extractor):
    """Plain text / markdown / json / yaml / xml — read bytes as UTF-8."""

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        raw = path.read_bytes()
        try:
            content = raw.decode("utf-8")
        except UnicodeDecodeError:
            content = raw.decode("utf-8", errors="replace")
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "bytes": len(raw)},
            )
        ]


@extractor("pdf", ["application/pdf"])
class PDFExtractor(Extractor):
    """PDF → text via pdfplumber (page-separated)."""

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        try:
            import pdfplumber  # type: ignore
        except ImportError as e:
            raise RuntimeError(
                "PDF extraction requires pdfplumber. Install with: pip install pdfplumber"
            ) from e

        pages_text: List[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                txt = page.extract_text() or ""
                pages_text.append(f"## Page {i}\n{txt}".strip())

        content = "\n\n".join(pages_text)
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "pages": len(pages_text)},
            )
        ]


@extractor("docx", [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
])
class DocxExtractor(Extractor):
    """DOCX/DOC → text via python-docx."""

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        try:
            import docx  # type: ignore
        except ImportError as e:
            raise RuntimeError(
                "DOCX extraction requires python-docx. Install with: pip install python-docx"
            ) from e

        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        content = "\n\n".join(paragraphs)
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "paragraphs": len(paragraphs)},
            )
        ]


@extractor("pptx", [
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
])
class PptxExtractor(Extractor):
    """PPTX/PPT → text via python-pptx (slide-separated)."""

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        try:
            import pptx  # type: ignore
        except ImportError as e:
            raise RuntimeError(
                "PPTX extraction requires python-pptx. Install with: pip install python-pptx"
            ) from e

        prs = pptx.Presentation(str(path))
        slides_text: List[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            texts.append(t)
            if texts:
                slides_text.append(f"## Slide {i}\n" + "\n".join(texts))

        content = "\n\n".join(slides_text)
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "slides": len(slides_text)},
            )
        ]


# ---------------------------------------------------------------------------
# Multimodal extractors (require model_caller)
# ---------------------------------------------------------------------------


@extractor("image", ["image/*"])
class ImageExtractor(Extractor):
    """Image → caption verbatim via a multimodal LLM.

    Calls `model_caller(model, prompt, images=[path])`. Falls back to an
    OCR-only verbatim (empty caption, marked deprecated) if no model_caller.
    """

    PROMPT = (
        "请详细描述这张图片的内容。包括：可见的文字（OCR）、主体对象、场景、"
        "图表数据、以及任何对知识库有用的元信息。输出纯文本，不要 markdown 标题。"
    )

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        date_iso, mtime = _file_mtime_iso(path)
        if not model or not model_caller:
            logger.warning(
                "ImageExtractor: no model/model_caller supplied for %s; "
                "writing empty-caption verbatim",
                path.name,
            )
            caption = ""
        else:
            try:
                caption = await model_caller(model, self.PROMPT, images=[path])
            except Exception as e:
                logger.exception("ImageExtractor model call failed for %s", path.name)
                raise RuntimeError(f"Image caption model call failed: {e}") from e

        content = (
            f"[Image: {path.name}]\n\n"
            f"**MIME**: {mime}\n\n"
            f"**Caption**:\n{caption}"
        )
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "model": model, "caption_len": len(caption)},
            )
        ]


@extractor("audio", ["audio/*"])
class AudioExtractor(Extractor):
    """Audio → transcription verbatim via an ASR/multimodal model.

    Calls `model_caller(model, prompt, images=[path])` — the caller is
    responsible for routing audio files to a model that accepts them
    (e.g., qwen-audio). If no model_caller is supplied, raises a clear error.
    """

    PROMPT = "请转录这段音频的内容。保留原文，不要总结。如果有多人对话，标注说话人。"

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        date_iso, mtime = _file_mtime_iso(path)
        if not model or not model_caller:
            raise RuntimeError(
                f"Audio extraction for {path.name} requires a multimodal/ASR model. "
                "Configure a multimodal_model on the space or pass a model_override."
            )
        try:
            transcript = await model_caller(model, self.PROMPT, images=[path])
        except Exception as e:
            logger.exception("AudioExtractor model call failed for %s", path.name)
            raise RuntimeError(f"Audio transcription model call failed: {e}") from e

        content = (
            f"[Audio: {path.name}]\n\n"
            f"**MIME**: {mime}\n\n"
            f"**Transcript**:\n{transcript}"
        )
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "model": model},
            )
        ]


# ---------------------------------------------------------------------------
# Video
# ---------------------------------------------------------------------------


async def _run_cmd(cmd: List[str]) -> int:
    """Run a subprocess quietly; returns the exit code."""
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.DEVNULL,
        stderr=asyncio.subprocess.DEVNULL,
    )
    return await proc.wait()


@extractor("video", ["video/*"])
class VideoExtractor(Extractor):
    """Video → 双策略:原生视频模型直连,不支持/失败则降级 ffmpeg 抽帧管线。

    - 原生直连:model 命中 NATIVE_MODEL_PATTERN(gemini / qwen-vl / qwen-omni)
      且文件 ≤ MAX_NATIVE_BYTES 时,经 model_caller(videos=[path]) 直接理解;
    - 抽帧管线:ffmpeg 按间隔抽帧(cap MAX_FRAMES)→ 分批走 image 通道逐批
      描述;音轨抽出转写(best-effort);合并为带片段结构的 verbatim。

    环境依赖:抽帧管线需要 ffmpeg 在 PATH 上。
    """

    NATIVE_MODEL_PATTERN = re.compile(r"gemini|qwen.*(vl|omni)", re.IGNORECASE)
    MAX_NATIVE_BYTES = 20 * 1024 * 1024
    MAX_FRAMES = 24
    FRAME_BATCH = 8
    FRAME_INTERVAL_SEC = 5

    PROMPT_NATIVE = (
        "请详细描述这个视频的内容:场景、事件发展过程、可见的文字、"
        "以及任何对知识库有用的关键信息。输出纯文本,不要 markdown 标题。"
    )
    PROMPT_FRAMES = (
        "这是同一个视频按时间顺序抽取的一组帧。请描述这组帧里发生的事情、"
        "场景变化、可见的文字与关键信息。输出纯文本,不要 markdown 标题。"
    )
    PROMPT_AUDIO = "请转录这段音频的内容。保留原文,不要总结。如果有多人对话,标注说话人。"

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        if not model or not model_caller:
            raise RuntimeError(
                f"Video extraction for {path.name} requires a multimodal model. "
                "Configure a multimodal_model on the space or pass a model_override."
            )
        date_iso, mtime = _file_mtime_iso(path)

        # 1. 原生视频模型直连(限小文件;失败/空输出降级抽帧)
        try:
            size = path.stat().st_size
        except OSError:
            size = 0
        if self.NATIVE_MODEL_PATTERN.search(model) and size <= self.MAX_NATIVE_BYTES:
            try:
                text = await model_caller(model, self.PROMPT_NATIVE, videos=[path])
                if text and text.strip():
                    content = (
                        f"[Video: {path.name}]\n\n"
                        f"**MIME**: {mime}\n\n"
                        f"**Content**:\n{text}"
                    )
                    return [
                        VerbatimSpec(
                            content=content,
                            source_file=path.name,
                            extract_mode=ExtractMode.UPLOAD,
                            content_date=date_iso,
                            source_mtime=mtime,
                            meta={"mime": mime, "model": model, "strategy": "native"},
                        )
                    ]
                logger.warning(
                    "VideoExtractor: native call returned empty for %s; "
                    "falling back to frame pipeline",
                    path.name,
                )
            except Exception as e:  # noqa: BLE001
                logger.warning(
                    "VideoExtractor: native call failed for %s (%s); "
                    "falling back to frame pipeline",
                    path.name,
                    e,
                )

        # 2. ffmpeg 抽帧 + 音轨管线
        if not shutil.which("ffmpeg"):
            raise RuntimeError(
                "Video extraction requires ffmpeg (not found on PATH), "
                "or a native video model (gemini / qwen-vl)."
            )
        import tempfile

        with tempfile.TemporaryDirectory(prefix="ks_video_") as tmp:
            tmp_dir = Path(tmp)
            frames = await self._extract_frames(path, tmp_dir)
            segments: List[str] = []
            for i in range(0, len(frames), self.FRAME_BATCH):
                batch = frames[i : i + self.FRAME_BATCH]
                desc = await model_caller(model, self.PROMPT_FRAMES, images=batch)
                segments.append(f"[片段 {i // self.FRAME_BATCH + 1}]\n{desc}")

            transcript = ""
            audio_path = await self._extract_audio(path, tmp_dir)
            if audio_path is not None:
                try:
                    transcript = await model_caller(
                        model, self.PROMPT_AUDIO, images=[audio_path]
                    )
                except Exception as e:  # noqa: BLE001
                    logger.warning(
                        "VideoExtractor: audio transcription failed for %s: %s",
                        path.name,
                        e,
                    )

            if not segments and not transcript.strip():
                raise RuntimeError(
                    f"Video extraction for {path.name} produced no content "
                    "(no frames, no audio track)"
                )
            parts = [f"[Video: {path.name}]\n\n**MIME**: {mime}"]
            if segments:
                parts.append("**画面**:\n" + "\n\n".join(segments))
            if transcript.strip():
                parts.append("**音轨转写**:\n" + transcript)
            return [
                VerbatimSpec(
                    content="\n\n".join(parts),
                    source_file=path.name,
                    extract_mode=ExtractMode.UPLOAD,
                    content_date=date_iso,
                    source_mtime=mtime,
                    meta={
                        "mime": mime,
                        "model": model,
                        "strategy": "frames",
                        "frames": len(frames),
                        "audio": bool(transcript.strip()),
                    },
                )
            ]

    async def _extract_frames(self, path: Path, out_dir: Path) -> List[Path]:
        """ffmpeg 按固定间隔抽帧(缩放至宽 768),上限 MAX_FRAMES。"""
        pattern = out_dir / "frame_%03d.jpg"
        rc = await _run_cmd(
            [
                "ffmpeg",
                "-y",
                "-i",
                str(path),
                "-vf",
                f"fps=1/{self.FRAME_INTERVAL_SEC},scale=768:-2",
                "-frames:v",
                str(self.MAX_FRAMES),
                str(pattern),
            ]
        )
        if rc != 0:
            logger.warning("VideoExtractor: ffmpeg frame extraction failed for %s", path.name)
            return []
        return sorted(out_dir.glob("frame_*.jpg"))

    async def _extract_audio(self, path: Path, out_dir: Path) -> Optional[Path]:
        """ffmpeg 抽出音轨为 16k 单声道 wav;无音轨返回 None。"""
        out = out_dir / "audio.wav"
        rc = await _run_cmd(
            ["ffmpeg", "-y", "-i", str(path), "-vn", "-ac", "1", "-ar", "16000", str(out)]
        )
        if rc != 0 or not out.exists() or out.stat().st_size == 0:
            return None
        return out
